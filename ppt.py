from pptx import Presentation
from pptx.util import Inches, Pt

def create_midterm_ppt():
    prs = Presentation()
    
    # 1. 封面
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "SDR无线电接收机设计与实现"
    subtitle.text = "2026届本科毕业设计中期答辩\n\n汇报人：贾晓晨 (22211428)\n专业：通信工程\n学院：电子信息工程学院，北京交通大学"

    # 2. 项目背景与核心任务
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "项目背景与核心任务"
    tf = slide.placeholders[1].text_frame
    tf.text = "项目背景："
    p = tf.add_paragraph()
    p.text = "SDR技术利用处理器通过软件编程实现调制、解调和滤波，提高灵活性并降低成本。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "核心任务：基于树莓派和RTL-SDR开发独立无线电接收机。"
    p = tf.add_paragraph()
    p.text = "关键技术指标："
    p = tf.add_paragraph()
    p.text = "频率：50-470MHz；解调：WFM, NFM, AM"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "支持直接输入(VF)及频道存储调用(VM，≥100个)"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "最大扫频宽度2MHz，音频功率≥0.5W，支持3.5mm及电池/USB供电"
    p.level = 1

    # 3. 总体方案与技术路线
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "总体方案与技术路线"
    tf = slide.placeholders[1].text_frame
    tf.text = "硬件架构：RTL-SDR采集射频，树莓派(Linux)主控。"
    p = tf.add_paragraph()
    p.text = "软件架构："
    p = tf.add_paragraph()
    p.text = "基带处理：GNU Radio构建信号流图完成下变频与解调。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "通信控制：ZMQ传输底层频谱数据，XMLRPC控制频率/增益。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "前端交互：Python + PyQtGraph开发亮色主题实时渲染UI。"
    p.level = 1

    # 4. 已完成的工作
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "已完成的工作"
    tf = slide.placeholders[1].text_frame
    tf.text = "软件与算法实现："
    p = tf.add_paragraph()
    p.text = "GNU Radio端成功验证WFM、NFM、AM解调流图。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "完成定制化UI开发，实现频谱图实时显示与自上而下的瀑布图。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "代码深度重构，将Python主程序拆分为多文件模块，提高可维护性。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "结构设计："
    p = tf.add_paragraph()
    p.text = "使用SolidWorks完成外壳建模，并为3D打印拆分为上下壳。"
    p.level = 1

    # 5. 攻克的技术难点
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "攻克的技术难点"
    tf = slide.placeholders[1].text_frame
    tf.text = "难点1：频谱数据对齐偏差"
    p = tf.add_paragraph()
    p.text = "确认GNU Radio输出的FFT已平移(Shifted)，移除Python端冗余的二次平移，精准还原频谱对齐。"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "难点2：高频UI刷新导致的CPU瓶颈"
    p = tf.add_paragraph()
    p.text = "引入PyQtGraph与ZMQ异步数据接收，大幅降低树莓派负载，保障渲染流畅度。"
    p.level = 1

    # 6. 尚需完成的工作及计划
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "尚需完成的工作及计划安排"
    tf = slide.placeholders[1].text_frame
    tf.text = "4月上旬：完成3D打印机壳制作与硬件实物组装。"
    p = tf.add_paragraph()
    p.text = "4月中旬：开发调制方式自动识别算法，实现100个频道的存储(VM)功能。"
    p = tf.add_paragraph()
    p.text = "4月下旬：调试0.5W音频放大电路，开展室外实地扫频与解调演示。"
    p = tf.add_paragraph()
    p.text = "5月份：整理软硬件测试图纸及数据，撰写毕业设计论文。"

    prs.save('SDR_Midterm_Presentation.pptx')
    print("PPT文件 'SDR_Midterm_Presentation.pptx' 生成成功！")

if __name__ == '__main__':
    create_midterm_ppt()